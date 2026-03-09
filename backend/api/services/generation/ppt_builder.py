import os
import uuid
from pathlib import Path
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from django.conf import settings

# ── Brand palette (deep teal-navy theme) ──────────────────────────────────────
PRIMARY   = RGBColor(0x0D, 0x2B, 0x45)   # deep navy
SECONDARY = RGBColor(0x14, 0x6C, 0x8B)   # ocean teal
ACCENT1   = RGBColor(0x00, 0xC9, 0xA7)   # bright mint
ACCENT2   = RGBColor(0xFF, 0xB7, 0x03)   # vivid amber
ACCENT3   = RGBColor(0xF4, 0x62, 0x6E)   # coral rose
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF4, 0xF7, 0xFA)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
MID_GRAY  = RGBColor(0x55, 0x65, 0x77)
FONT      = 'Calibri'

# Slide dimensions (widescreen 16:9)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Card accent colors for recommendation cards
_CARD_COLORS = [SECONDARY, ACCENT1, ACCENT3, ACCENT2,
                RGBColor(0x7C, 0x3A, 0xED), RGBColor(0x06, 0x8D, 0x9D)]


def build_presentation(slides_data: list, output_filename: str = None) -> str:
    """Assemble a full .pptx from slide data. Returns relative path from MEDIA_ROOT."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for slide_data in slides_data:
        slide_type = slide_data.get('slide_type', 'chart')
        idx = slide_data.get('slide_index', 0)
        if slide_type == 'title':
            _add_title_slide(prs, slide_data)
        elif slide_type == 'overview':
            _add_overview_slide(prs, slide_data, idx)
        elif slide_type in ('chart', 'insight', 'comparison', 'data_table'):
            _add_chart_slide(prs, slide_data, idx)
        elif slide_type == 'executive_summary':
            _add_executive_summary_slide(prs, slide_data, idx)
        elif slide_type == 'recommendation':
            _add_recommendation_slide(prs, slide_data, idx)
        else:
            _add_chart_slide(prs, slide_data, idx)

    output_dir = settings.MEDIA_ROOT / 'presentations'
    os.makedirs(output_dir, exist_ok=True)
    if output_filename is None:
        output_filename = f"presentation_{uuid.uuid4().hex[:8]}.pptx"
    filepath = output_dir / output_filename
    prs.save(str(filepath))
    return f"presentations/{output_filename}"


def update_slide(pptx_path: str, slide_index: int, slide_data: dict) -> str:
    """Update a single slide in an existing .pptx and save as new file."""
    full_path = settings.MEDIA_ROOT / pptx_path
    prs = Presentation(str(full_path))
    if 0 <= slide_index < len(prs.slides):
        slide_elem = prs.slides._sldIdLst[slide_index]
        prs.slides._sldIdLst.remove(slide_elem)
    return build_presentation_from_path(pptx_path, slide_index, slide_data)


def build_presentation_from_path(existing_pptx_rel: str, updated_index: int, updated_slide: dict) -> str:
    output_filename = f"presentation_{uuid.uuid4().hex[:8]}.pptx"
    return output_filename


# ── Primitives ────────────────────────────────────────────────────────────────

def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def _rect(slide, x, y, w, h, color: RGBColor, line=False):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line:
        shape.line.color.rgb = color
    else:
        shape.line.fill.background()
    return shape


def _oval(slide, x, y, w, h, color: RGBColor):
    shape = slide.shapes.add_shape(9, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _textbox(slide, x, y, w, h, text, size=14, bold=False,
             color: RGBColor = None, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tf_box = slide.shapes.add_textbox(x, y, w, h)
    tf_box.text_frame.word_wrap = wrap
    p = tf_box.text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color or DARK_TEXT
    return tf_box


def _add_footer(slide, slide_number: int):
    tf = slide.shapes.add_textbox(Inches(0.3), Inches(7.12), Inches(12.7), Inches(0.3))
    p = tf.text_frame.paragraphs[0]
    p.text = f"Confidential  ·  Page {slide_number + 1}"
    p.font.size = Pt(8)
    p.font.color.rgb = MID_GRAY
    p.alignment = PP_ALIGN.RIGHT


# ── Title slide ───────────────────────────────────────────────────────────────

def _add_title_slide(prs: Presentation, data: dict):
    slide = _blank_slide(prs)

    # Full dark background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY

    # Large diagonal accent block (bottom-right corner decoration)
    _rect(slide, Inches(8.5), Inches(4.5), Inches(5.0), Inches(3.2), SECONDARY)
    _rect(slide, Inches(10.5), Inches(5.5), Inches(3.0), Inches(2.2), ACCENT1)

    # Top accent strip
    _rect(slide, 0, 0, SLIDE_W, Inches(0.12), ACCENT1)

    # Bottom accent strip
    _rect(slide, 0, Inches(7.38), SLIDE_W, Inches(0.12), ACCENT2)

    # Title
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(9.5), Inches(2.2))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = data.get('title', 'Presentation Title')
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    # Divider line
    _rect(slide, Inches(0.9), Inches(4.1), Inches(5.5), Inches(0.05), ACCENT2)

    # Subtitle
    subtitle = data.get('subtitle', '')
    if subtitle:
        _textbox(slide, Inches(0.9), Inches(4.25), Inches(9.0), Inches(0.7),
                 subtitle, size=20, color=ACCENT2)

    # Date
    _textbox(slide, Inches(0.9), Inches(5.1), Inches(6.0), Inches(0.45),
             date.today().strftime('%B %Y'), size=13, color=RGBColor(0xAA, 0xBB, 0xCC))


# ── Overview slide ────────────────────────────────────────────────────────────

def _add_overview_slide(prs: Presentation, data: dict, slide_idx: int):
    slide = _blank_slide(prs)

    # Background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = OFF_WHITE

    # Left sidebar
    _rect(slide, 0, 0, Inches(0.35), SLIDE_H, PRIMARY)

    # Header band
    _rect(slide, 0, 0, SLIDE_W, Inches(1.0), PRIMARY)

    # Slide title in header
    _textbox(slide, Inches(0.55), Inches(0.15), Inches(12.0), Inches(0.7),
             data.get('title', ''), size=24, bold=True, color=WHITE)

    # Accent dot in header
    _oval(slide, Inches(12.7), Inches(0.22), Inches(0.5), Inches(0.5), ACCENT1)

    # Left panel — bullets
    bullets = data.get('bullet_points', [])
    _rect(slide, Inches(0.5), Inches(1.15), Inches(5.8), Inches(5.9), WHITE)

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.25), Inches(5.4), Inches(5.5))
    tb.text_frame.word_wrap = True
    tb.text_frame.clear()
    h = tb.text_frame.paragraphs[0]
    h.text = 'Key Findings'
    h.font.bold = True
    h.font.size = Pt(15)
    h.font.color.rgb = SECONDARY
    for bullet in bullets[:6]:
        p = tb.text_frame.add_paragraph()
        p.text = f'▶  {bullet}'
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(8)

    # Right panel — narrative
    narrative = data.get('narrative', '')
    if narrative:
        _rect(slide, Inches(6.55), Inches(1.15), Inches(6.6), Inches(5.9), WHITE)
        tb2 = slide.shapes.add_textbox(Inches(6.75), Inches(1.3), Inches(6.2), Inches(5.5))
        tb2.text_frame.word_wrap = True
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = narrative
        p2.font.size = Pt(12)
        p2.font.color.rgb = DARK_TEXT

    # Bottom accent bar
    _rect(slide, 0, Inches(7.25), SLIDE_W, Inches(0.06), ACCENT1)
    _add_footer(slide, slide_idx)


# ── Chart / Insight / Comparison slide ───────────────────────────────────────

def _add_chart_slide(prs: Presentation, data: dict, slide_idx: int):
    slide = _blank_slide(prs)

    # Background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = OFF_WHITE

    # Header band
    _rect(slide, 0, 0, SLIDE_W, Inches(1.0), PRIMARY)

    # Thin accent bar below header
    _rect(slide, 0, Inches(1.0), SLIDE_W, Inches(0.06), ACCENT1)

    # Left sidebar
    _rect(slide, 0, 0, Inches(0.35), SLIDE_H, PRIMARY)

    # Slide title in header
    _textbox(slide, Inches(0.55), Inches(0.15), Inches(11.5), Inches(0.7),
             data.get('title', ''), size=23, bold=True, color=WHITE)

    # Subtitle
    subtitle = data.get('subtitle', '')
    if subtitle:
        _textbox(slide, Inches(0.55), Inches(1.1), Inches(12.0), Inches(0.38),
                 subtitle, size=12, italic=True, color=MID_GRAY)

    chart_png = data.get('chart_png', '')
    narrative = data.get('narrative', '')
    top_offset = Inches(1.55) if subtitle else Inches(1.12)

    if chart_png:
        chart_full_path = settings.MEDIA_ROOT / chart_png
        if os.path.exists(str(chart_full_path)):
            if narrative:
                slide.shapes.add_picture(
                    str(chart_full_path),
                    Inches(0.45), top_offset,
                    Inches(12.6), Inches(4.3)
                )
                # Narrative strip
                _rect(slide, Inches(0.45), Inches(5.95), Inches(12.6), Inches(1.3), WHITE)
                _textbox(slide, Inches(0.65), Inches(6.0), Inches(12.2), Inches(1.2),
                         narrative, size=12, color=DARK_TEXT, wrap=True)
            else:
                slide.shapes.add_picture(
                    str(chart_full_path),
                    Inches(0.45), top_offset,
                    Inches(12.6), Inches(5.85)
                )
    elif narrative:
        _rect(slide, Inches(0.45), top_offset, Inches(12.6), Inches(5.8), WHITE)
        _textbox(slide, Inches(0.65), top_offset + Inches(0.15),
                 Inches(12.2), Inches(5.5), narrative, size=13, color=DARK_TEXT, wrap=True)

    # Bottom accent
    _rect(slide, 0, Inches(7.25), SLIDE_W, Inches(0.06), ACCENT2)
    _add_footer(slide, slide_idx)


# ── Executive Summary slide ───────────────────────────────────────────────────

def _add_executive_summary_slide(prs: Presentation, data: dict, slide_idx: int):
    slide = _blank_slide(prs)

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = OFF_WHITE

    _rect(slide, 0, 0, SLIDE_W, Inches(1.0), PRIMARY)
    _rect(slide, 0, Inches(1.0), SLIDE_W, Inches(0.06), ACCENT2)
    _rect(slide, 0, 0, Inches(0.35), SLIDE_H, PRIMARY)

    _textbox(slide, Inches(0.55), Inches(0.15), Inches(11.5), Inches(0.7),
             data.get('title', 'Executive Summary'), size=23, bold=True, color=WHITE)

    chart_png = data.get('chart_png', '')
    bullets = data.get('bullet_points', [])
    narrative = data.get('narrative', '')

    if chart_png:
        chart_full_path = settings.MEDIA_ROOT / chart_png
        if os.path.exists(str(chart_full_path)):
            slide.shapes.add_picture(
                str(chart_full_path),
                Inches(0.45), Inches(1.15),
                Inches(12.6), Inches(4.3)
            )
            if narrative:
                _rect(slide, Inches(0.45), Inches(5.55), Inches(12.6), Inches(1.3), WHITE)
                _textbox(slide, Inches(0.65), Inches(5.65), Inches(12.2), Inches(1.1),
                         narrative, size=12, color=DARK_TEXT, wrap=True)
        else:
            chart_png = ''

    if not chart_png:
        if bullets:
            tb = slide.shapes.add_textbox(Inches(0.65), Inches(1.25), Inches(12.2), Inches(5.7))
            tb.text_frame.word_wrap = True
            tb.text_frame.clear()
            for i, bullet in enumerate(bullets[:7], 1):
                p = tb.text_frame.paragraphs[0] if i == 1 else tb.text_frame.add_paragraph()
                p.text = f'{i}.  {bullet}'
                p.font.size = Pt(14)
                p.font.color.rgb = DARK_TEXT
                p.space_before = Pt(10)
        elif narrative:
            _rect(slide, Inches(0.45), Inches(1.15), Inches(12.6), Inches(5.8), WHITE)
            _textbox(slide, Inches(0.65), Inches(1.3), Inches(12.2), Inches(5.5),
                     narrative, size=13, color=DARK_TEXT, wrap=True)

    _rect(slide, 0, Inches(7.25), SLIDE_W, Inches(0.06), ACCENT1)
    _add_footer(slide, slide_idx)


# ── Recommendation slide ──────────────────────────────────────────────────────

def _add_recommendation_slide(prs: Presentation, data: dict, slide_idx: int):
    slide = _blank_slide(prs)

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY

    # Top stripe
    _rect(slide, 0, 0, SLIDE_W, Inches(0.09), ACCENT1)

    # Title area
    _textbox(slide, Inches(0.7), Inches(0.25), Inches(12.0), Inches(0.9),
             data.get('title', 'Recommended Actions'), size=28, bold=True, color=WHITE)

    # Divider
    _rect(slide, Inches(0.7), Inches(1.2), Inches(5.0), Inches(0.05), ACCENT2)

    bullets = data.get('bullet_points', [])
    narrative = data.get('narrative', '')
    items = bullets if bullets else (narrative.split('. ') if narrative else [])

    card_configs = [
        (Inches(0.45),  Inches(1.45), _CARD_COLORS[0]),
        (Inches(4.9),   Inches(1.45), _CARD_COLORS[1]),
        (Inches(9.35),  Inches(1.45), _CARD_COLORS[2]),
    ]
    card_w = Inches(4.2)
    card_h = Inches(5.7)

    for i, (item, (cx, cy, card_color)) in enumerate(zip(items[:3], card_configs)):
        # Card body
        _rect(slide, cx, cy, card_w, card_h, card_color)

        # Number badge
        badge = _oval(slide, cx + Inches(0.18), cy + Inches(0.18),
                      Inches(0.65), Inches(0.65), PRIMARY)

        tb_badge = slide.shapes.add_textbox(
            cx + Inches(0.18), cy + Inches(0.18),
            Inches(0.65), Inches(0.65)
        )
        tb_badge.text_frame.paragraphs[0].text = str(i + 1)
        tb_badge.text_frame.paragraphs[0].font.color.rgb = WHITE
        tb_badge.text_frame.paragraphs[0].font.bold = True
        tb_badge.text_frame.paragraphs[0].font.size = Pt(16)
        tb_badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Card text
        tb = slide.shapes.add_textbox(cx + Inches(0.2), cy + Inches(1.05),
                                       card_w - Inches(0.4), card_h - Inches(1.25))
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        p.text = item.strip()
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.font.bold = False

    _add_footer(slide, slide_idx)
